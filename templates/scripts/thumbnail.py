#!/usr/bin/env python3
"""Create thumbnail grids from PowerPoint presentation slides.

Usage:
    python scripts/thumbnail.py input.pptx [output_prefix] [--cols N]

Creates a grid layout of slide thumbnails for quick visual analysis.
Labels each thumbnail with its XML filename.

Requires: Pillow, python-pptx
Optional: LibreOffice (soffice) + pdftoppm (poppler) for higher-fidelity rendering
"""

import argparse
import os
import platform
import shutil
import subprocess
import sys
import tempfile
import zipfile
from io import BytesIO
from pathlib import Path
from xml.dom import minidom

from PIL import Image, ImageDraw, ImageFont

THUMBNAIL_WIDTH = 300
RENDER_WIDTH_PPTX = 1920
CONVERSION_DPI = 100
MAX_COLS = 6
DEFAULT_COLS = 3
JPEG_QUALITY = 95
GRID_PADDING = 20
BORDER_WIDTH = 2
FONT_SIZE_RATIO = 0.10
LABEL_PADDING_RATIO = 0.4


def get_soffice_path() -> str:
    """Find soffice binary for the current platform."""
    system = platform.system()

    if system == "Darwin":
        candidates = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            os.path.expanduser("~/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        ]
    elif system == "Windows":
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
    else:  # Linux / WSL
        candidates = [
            "/usr/bin/soffice",
            "/usr/lib/libreoffice/program/soffice",
        ]

    for p in candidates:
        if os.path.isfile(p):
            return p

    # Fallback to PATH
    return "soffice"


def get_slide_info(pptx_path: Path) -> list[dict]:
    """Get ordered slide info from the presentation."""
    with zipfile.ZipFile(pptx_path, "r") as zf:
        rels_content = zf.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        rels_dom = minidom.parseString(rels_content)

        rid_to_slide = {}
        for rel in rels_dom.getElementsByTagName("Relationship"):
            rid = rel.getAttribute("Id")
            target = rel.getAttribute("Target")
            rel_type = rel.getAttribute("Type")
            if "slide" in rel_type and target.startswith("slides/"):
                rid_to_slide[rid] = target.replace("slides/", "")

        pres_content = zf.read("ppt/presentation.xml").decode("utf-8")
        pres_dom = minidom.parseString(pres_content)

        slides = []
        for sld_id in pres_dom.getElementsByTagName("p:sldId"):
            rid = sld_id.getAttribute("r:id")
            if rid in rid_to_slide:
                slides.append({"name": rid_to_slide[rid]})

        return slides


def has_soffice() -> bool:
    """Check if LibreOffice is available."""
    soffice = get_soffice_path()
    if soffice != "soffice":
        return True
    return shutil.which("soffice") is not None


def convert_to_images_soffice(pptx_path: Path, temp_dir: Path) -> list[Path]:
    """Convert .pptx to slide images via LibreOffice + pdftoppm."""
    soffice = get_soffice_path()
    pdf_path = temp_dir / f"{pptx_path.stem}.pdf"

    result = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(temp_dir), str(pptx_path)],
        capture_output=True, text=True,
    )
    if result.returncode != 0 or not pdf_path.exists():
        return []

    if not shutil.which("pdftoppm"):
        return []

    result = subprocess.run(
        ["pdftoppm", "-jpeg", "-r", str(CONVERSION_DPI), str(pdf_path), str(temp_dir / "slide")],
        capture_output=True, text=True,
    )
    if result.returncode != 0:
        return []

    return sorted(temp_dir.glob("slide-*.jpg"))


def convert_to_images_pptx(pptx_path: Path, temp_dir: Path) -> list[Path]:
    """Enhanced slide renderer using python-pptx + Pillow.

    Renders backgrounds (with layout/master fallback), embedded images,
    auto shapes (rectangles, ovals, rounded rects) with fills and borders,
    tables with cell styling, and text with actual font sizes, colors, bold,
    and alignment. Uses system TrueType fonts when available.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width or Emu(9144000)   # default 10"
    slide_h = prs.slide_height or Emu(6858000)  # default 7.5"

    render_w = RENDER_WIDTH_PPTX
    render_h = int(render_w * (slide_h / slide_w))

    images = []
    for i, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (render_w, render_h), "#FFFFFF")
        draw = ImageDraw.Draw(img)

        # 1. Background (slide -> layout -> master fallback)
        bg_color = _render_slide_bg(slide, img, draw, render_w, render_h)

        # 2. Shapes in z-order
        for shape in slide.shapes:
            _render_shape(shape, img, draw, slide_w, slide_h,
                          render_w, render_h, bg_color)

        out_path = temp_dir / f"slide-{i:02d}.jpg"
        img.save(str(out_path), quality=JPEG_QUALITY)
        images.append(out_path)

    return images


# ---- Font resolution ---------------------------------------------------------

_font_cache: dict = {}
_font_dirs_cache: list[Path] | None = None


def _get_font_dirs() -> list[Path]:
    """Return system font directories for the current platform."""
    global _font_dirs_cache
    if _font_dirs_cache is not None:
        return _font_dirs_cache
    candidates: list[Path] = []
    system = platform.system()
    if system == "Darwin":
        candidates = [
            Path("/System/Library/Fonts/Supplemental"),
            Path("/Library/Fonts"),
            Path.home() / "Library" / "Fonts",
        ]
    elif system == "Windows":
        windir = os.environ.get("WINDIR", r"C:\Windows")
        candidates = [Path(windir) / "Fonts"]
    else:
        candidates = [
            Path("/usr/share/fonts/truetype"),
            Path("/usr/share/fonts"),
            Path("/usr/local/share/fonts"),
            Path.home() / ".local" / "share" / "fonts",
        ]
    _font_dirs_cache = [d for d in candidates if d.is_dir()]
    return _font_dirs_cache


def _resolve_font(name: str | None, size_pt: float,
                  bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """Load a system TrueType font by name; fall back to Pillow built-in."""
    size_px = max(int(size_pt * 4 / 3), 8)
    key = (name, size_px, bold)
    if key in _font_cache:
        return _font_cache[key]

    if name:
        clean = name.replace(" ", "")
        names = [name, clean]
        suffixes = (["-Bold", "Bold", " Bold", "bd", "-bold"]
                    if bold else ["", "-Regular", "Regular"])
        for font_dir in _get_font_dirs():
            for n in names:
                for suf in suffixes:
                    for ext in (".ttf", ".ttc", ".otf"):
                        p = font_dir / f"{n}{suf}{ext}"
                        if p.exists():
                            try:
                                font = ImageFont.truetype(str(p), size_px)
                                _font_cache[key] = font
                                return font
                            except Exception:
                                continue

    try:
        font = ImageFont.load_default(size=size_px)
    except TypeError:
        font = ImageFont.load_default()
    _font_cache[key] = font
    return font


# ---- Coordinate / color helpers ---------------------------------------------

def _emu_to_px(emu: int | None, slide_dim: int, render_dim: int) -> int:
    """Convert an EMU measurement to pixels."""
    if not emu or slide_dim == 0:
        return 0
    return int(emu / slide_dim * render_dim)


def _extract_fill_color(fill_obj) -> str | None:
    """Extract dominant color from a python-pptx FillFormat."""
    try:
        if fill_obj is None or fill_obj.type is None:
            return None
        fc = fill_obj.fore_color
        if fc and fc.rgb:
            return f"#{fc.rgb}"
    except Exception:
        pass
    return None


def _extract_line_props(line_obj) -> tuple[str | None, int]:
    """Extract color and pixel width from a python-pptx LineFormat."""
    color, width_px = None, 0
    try:
        if line_obj is None:
            return None, 0
        if line_obj.fill and line_obj.fill.type is not None:
            if line_obj.color and line_obj.color.rgb:
                color = f"#{line_obj.color.rgb}"
        if line_obj.width:
            width_px = max(1, int(line_obj.width / 914400 * 96))  # EMU -> px
    except Exception:
        pass
    return color, width_px


def _is_light(hex_color: str) -> bool:
    """Check if a hex color is perceptually light."""
    c = hex_color.lstrip("#")
    if len(c) < 6:
        return True
    try:
        r, g, b = int(c[:2], 16), int(c[2:4], 16), int(c[4:], 16)
        return (r * 299 + g * 587 + b * 114) / 1000 > 128
    except ValueError:
        return True


def _wrap_text(draw: ImageDraw.Draw, text: str, font, max_width: int) -> list[str]:
    """Word-wrap text to fit within max_width pixels."""
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        test = f"{current} {word}".strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [text]


# ---- Background -------------------------------------------------------------

def _render_slide_bg(slide, img: Image.Image, draw: ImageDraw.Draw,
                     w: int, h: int) -> str | None:
    """Render slide background with layout/master fallback. Return color."""
    color = None
    for source in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            bg = source.background
            if bg and bg.fill and bg.fill.type is not None:
                c = bg.fill.fore_color
                if c and c.rgb:
                    color = f"#{c.rgb}"
                    break
        except Exception:
            continue
    if color:
        draw.rectangle([(0, 0), (w, h)], fill=color)
    return color


# ---- Shape renderers --------------------------------------------------------

def _render_picture(shape, img: Image.Image,
                    slide_w: int, slide_h: int,
                    render_w: int, render_h: int) -> None:
    """Render an embedded image at its correct position and size."""
    try:
        blob = shape.image.blob
        with Image.open(BytesIO(blob)) as pic:
            if pic.mode not in ("RGB", "RGBA"):
                pic = pic.convert("RGBA")
            x = _emu_to_px(shape.left, slide_w, render_w)
            y = _emu_to_px(shape.top, slide_h, render_h)
            w = _emu_to_px(shape.width, slide_w, render_w)
            h = _emu_to_px(shape.height, slide_h, render_h)
            if w > 0 and h > 0:
                pic = pic.resize((w, h), Image.Resampling.LANCZOS)
                if pic.mode == "RGBA":
                    img.paste(pic, (x, y), pic)
                else:
                    img.paste(pic, (x, y))
    except Exception:
        pass


def _draw_rounded_rect(draw: ImageDraw.Draw,
                       xy: tuple[int, int, int, int], radius: int,
                       fill=None, outline=None, width: int = 1) -> None:
    """Draw a rounded rectangle (uses Pillow 10+ native, manual fallback)."""
    x0, y0, x1, y1 = xy
    r = min(radius, (x1 - x0) // 2, (y1 - y0) // 2)
    if r <= 0:
        draw.rectangle([(x0, y0), (x1, y1)], fill=fill,
                       outline=outline, width=width)
        return

    # Pillow >= 10.0 has draw.rounded_rectangle
    try:
        draw.rounded_rectangle([(x0, y0), (x1, y1)], radius=r,
                               fill=fill, outline=outline, width=width)
        return
    except AttributeError:
        pass

    # Manual fallback
    if fill:
        draw.rectangle([(x0 + r, y0), (x1 - r, y1)], fill=fill)
        draw.rectangle([(x0, y0 + r), (x1, y1 - r)], fill=fill)
        draw.pieslice([(x0, y0), (x0 + 2*r, y0 + 2*r)], 180, 270, fill=fill)
        draw.pieslice([(x1 - 2*r, y0), (x1, y0 + 2*r)], 270, 360, fill=fill)
        draw.pieslice([(x0, y1 - 2*r), (x0 + 2*r, y1)], 90, 180, fill=fill)
        draw.pieslice([(x1 - 2*r, y1 - 2*r), (x1, y1)], 0, 90, fill=fill)
    if outline:
        draw.arc([(x0, y0), (x0 + 2*r, y0 + 2*r)], 180, 270,
                 fill=outline, width=width)
        draw.arc([(x1 - 2*r, y0), (x1, y0 + 2*r)], 270, 360,
                 fill=outline, width=width)
        draw.arc([(x0, y1 - 2*r), (x0 + 2*r, y1)], 90, 180,
                 fill=outline, width=width)
        draw.arc([(x1 - 2*r, y1 - 2*r), (x1, y1)], 0, 90,
                 fill=outline, width=width)
        draw.line([(x0 + r, y0), (x1 - r, y0)], fill=outline, width=width)
        draw.line([(x0 + r, y1), (x1 - r, y1)], fill=outline, width=width)
        draw.line([(x0, y0 + r), (x0, y1 - r)], fill=outline, width=width)
        draw.line([(x1, y0 + r), (x1, y1 - r)], fill=outline, width=width)


def _render_auto_shape(shape, draw: ImageDraw.Draw,
                       slide_w: int, slide_h: int,
                       render_w: int, render_h: int) -> None:
    """Render an auto shape with proper geometry, fill, and border."""
    x = _emu_to_px(shape.left, slide_w, render_w)
    y = _emu_to_px(shape.top, slide_h, render_h)
    w = _emu_to_px(shape.width, slide_w, render_w)
    h = _emu_to_px(shape.height, slide_h, render_h)
    if w <= 0 or h <= 0:
        return

    fill_color = _extract_fill_color(shape.fill)
    line_color, line_width = _extract_line_props(shape.line)
    if not fill_color and not line_color:
        return  # invisible shape (no fill, no border — likely a text box)

    # Detect shape sub-type
    auto_type = None
    try:
        auto_type = shape.auto_shape_type
    except Exception:
        pass

    try:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        if auto_type == MSO_AUTO_SHAPE_TYPE.OVAL:
            draw.ellipse([(x, y), (x + w, y + h)],
                         fill=fill_color, outline=line_color,
                         width=line_width or 1)
        elif auto_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
            radius = min(w, h) // 5
            _draw_rounded_rect(draw, (x, y, x + w, y + h), radius,
                               fill=fill_color, outline=line_color,
                               width=line_width or 1)
        elif w <= 4 or h <= 4:
            # Thin shape -> accent bar / divider line
            draw.rectangle([(x, y), (x + w, y + h)],
                           fill=fill_color or line_color)
        else:
            draw.rectangle([(x, y), (x + w, y + h)],
                           fill=fill_color, outline=line_color,
                           width=line_width or 1)
    except Exception:
        if fill_color or line_color:
            draw.rectangle([(x, y), (x + w, y + h)],
                           fill=fill_color, outline=line_color,
                           width=line_width or 1)


def _render_table(shape, draw: ImageDraw.Draw,
                  slide_w: int, slide_h: int,
                  render_w: int, render_h: int) -> None:
    """Render a table with cell backgrounds, borders, and text."""
    try:
        table = shape.table
        x0 = _emu_to_px(shape.left, slide_w, render_w)
        y0 = _emu_to_px(shape.top, slide_h, render_h)

        col_widths = [_emu_to_px(c.width, slide_w, render_w)
                      for c in table.columns]
        row_heights = [_emu_to_px(r.height, slide_h, render_h)
                       for r in table.rows]

        cy = y0
        for r_idx, row in enumerate(table.rows):
            cx = x0
            rh = row_heights[r_idx] if r_idx < len(row_heights) else 30
            for c_idx, cell in enumerate(row.cells):
                cw = col_widths[c_idx] if c_idx < len(col_widths) else 80

                # Cell fill
                cell_fill = None
                try:
                    cell_fill = _extract_fill_color(cell.fill)
                except Exception:
                    pass

                draw.rectangle([(cx, cy), (cx + cw, cy + rh)],
                               fill=cell_fill, outline="#CCCCCC", width=1)

                # Cell text
                text = cell.text.strip()
                if text:
                    text_color = ("#FFFFFF" if cell_fill
                                  and not _is_light(cell_fill) else "#000000")
                    font_size, font_name, is_bold = 11, None, False
                    try:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    font_size = int(run.font.size.pt)
                                if run.font.name:
                                    font_name = run.font.name
                                if run.font.bold:
                                    is_bold = True
                                break
                            break
                    except Exception:
                        pass
                    font = _resolve_font(font_name, font_size, is_bold)
                    pad = 4
                    lines = _wrap_text(draw, text, font, cw - pad * 2)
                    for j, line in enumerate(lines[:3]):
                        draw.text((cx + pad, cy + pad + j * (font.size + 2)),
                                  line, fill=text_color, font=font)
                cx += cw
            cy += rh
    except Exception:
        pass


def _render_text_frame(shape, draw: ImageDraw.Draw,
                       slide_w: int, slide_h: int,
                       render_w: int, render_h: int,
                       bg_color: str | None) -> None:
    """Render a shape's text frame with paragraph-level formatting."""
    tf = shape.text_frame
    x = _emu_to_px(shape.left, slide_w, render_w)
    y = _emu_to_px(shape.top, slide_h, render_h)
    w = _emu_to_px(shape.width, slide_w, render_w)
    h = _emu_to_px(shape.height, slide_h, render_h)
    if w <= 0:
        return

    # Effective background for contrast detection
    shape_fill = _extract_fill_color(shape.fill)
    effective_bg = shape_fill or bg_color
    default_color = "#000000"
    if effective_bg and not _is_light(effective_bg):
        default_color = "#FFFFFF"

    padding = 6
    cursor_y = y + padding

    for para in tf.paragraphs:
        if cursor_y >= y + h:
            break

        text = para.text.strip()
        if not text:
            cursor_y += 10
            continue

        # Extract formatting from first run that has each property
        font_size, font_name, is_bold, text_color = 14, None, False, default_color
        for run in para.runs:
            if run.font.size:
                font_size = int(run.font.size.pt)
                break
        for run in para.runs:
            if run.font.name:
                font_name = run.font.name
                break
        for run in para.runs:
            if run.font.bold:
                is_bold = True
                break
        for run in para.runs:
            try:
                if run.font.color and run.font.color.rgb:
                    text_color = f"#{run.font.color.rgb}"
                    break
            except Exception:
                pass

        font = _resolve_font(font_name, font_size, is_bold)
        line_h = font.size + 4

        # Paragraph alignment
        align = "left"
        try:
            from pptx.enum.text import PP_ALIGN
            if para.alignment == PP_ALIGN.CENTER:
                align = "center"
            elif para.alignment == PP_ALIGN.RIGHT:
                align = "right"
        except Exception:
            pass

        lines = _wrap_text(draw, text, font, w - padding * 2)
        for line in lines:
            if cursor_y + line_h > y + h:
                break
            tx = x + padding
            if align == "center":
                bbox = draw.textbbox((0, 0), line, font=font)
                tx = x + (w - (bbox[2] - bbox[0])) // 2
            elif align == "right":
                bbox = draw.textbbox((0, 0), line, font=font)
                tx = x + w - (bbox[2] - bbox[0]) - padding
            draw.text((tx, cursor_y), line, fill=text_color, font=font)
            cursor_y += line_h

        cursor_y += 4  # paragraph spacing


# ---- Shape dispatcher -------------------------------------------------------

def _render_shape(shape, img: Image.Image, draw: ImageDraw.Draw,
                  slide_w: int, slide_h: int,
                  render_w: int, render_h: int,
                  bg_color: str | None = None) -> None:
    """Dispatch rendering based on shape type."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        st = shape.shape_type
    except Exception:
        return

    # Embedded images
    if st == MSO_SHAPE_TYPE.PICTURE:
        _render_picture(shape, img, slide_w, slide_h, render_w, render_h)
        return

    # Tables
    if st == MSO_SHAPE_TYPE.TABLE:
        _render_table(shape, draw, slide_w, slide_h, render_w, render_h)
        return

    # Groups — recurse into children
    if st == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                _render_shape(child, img, draw, slide_w, slide_h,
                              render_w, render_h, bg_color)
        except Exception:
            pass
        return

    # Auto shapes, text boxes, placeholders, freeforms
    # Step 1: shape fill and border
    _render_auto_shape(shape, draw, slide_w, slide_h, render_w, render_h)

    # Step 2: text on top
    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            _render_text_frame(shape, draw, slide_w, slide_h,
                               render_w, render_h, bg_color)
    except Exception:
        pass


def convert_to_images(pptx_path: Path, temp_dir: Path) -> list[Path]:
    """Convert .pptx to slide images. Uses LibreOffice if available, otherwise enhanced built-in renderer."""
    if has_soffice():
        images = convert_to_images_soffice(pptx_path, temp_dir)
        if images:
            return images
        print("NOTE: LibreOffice conversion failed, using built-in renderer.", file=sys.stderr)
    else:
        print("NOTE: Using enhanced built-in renderer (shapes, images, tables, text). Install LibreOffice + Poppler for full visual fidelity.", file=sys.stderr)

    return convert_to_images_pptx(pptx_path, temp_dir)


def create_grid(
    slides: list[tuple[Path, str]], cols: int, width: int
) -> Image.Image:
    """Create a thumbnail grid image."""
    font_size = int(width * FONT_SIZE_RATIO)
    label_padding = int(font_size * LABEL_PADDING_RATIO)

    with Image.open(slides[0][0]) as img:
        aspect = img.height / img.width
    height = int(width * aspect)

    rows = (len(slides) + cols - 1) // cols
    grid_w = cols * width + (cols + 1) * GRID_PADDING
    grid_h = rows * (height + font_size + label_padding * 2) + (rows + 1) * GRID_PADDING

    grid = Image.new("RGB", (grid_w, grid_h), "white")
    draw = ImageDraw.Draw(grid)

    try:
        font = ImageFont.load_default(size=font_size)
    except Exception:
        font = ImageFont.load_default()

    for i, (img_path, slide_name) in enumerate(slides):
        row, col = i // cols, i % cols
        x = col * width + (col + 1) * GRID_PADDING
        y_base = row * (height + font_size + label_padding * 2) + (row + 1) * GRID_PADDING

        # Draw label
        bbox = draw.textbbox((0, 0), slide_name, font=font)
        text_w = bbox[2] - bbox[0]
        draw.text(
            (x + (width - text_w) // 2, y_base + label_padding),
            slide_name, fill="black", font=font,
        )

        # Draw thumbnail
        y_thumbnail = y_base + label_padding + font_size + label_padding
        with Image.open(img_path) as img:
            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            w, h = img.size
            tx = x + (width - w) // 2
            ty = y_thumbnail + (height - h) // 2
            grid.paste(img, (tx, ty))

            if BORDER_WIDTH > 0:
                draw.rectangle(
                    [(tx - BORDER_WIDTH, ty - BORDER_WIDTH),
                     (tx + w + BORDER_WIDTH - 1, ty + h + BORDER_WIDTH - 1)],
                    outline="gray", width=BORDER_WIDTH,
                )

    return grid


def main():
    parser = argparse.ArgumentParser(description="Create thumbnail grids from PowerPoint slides.")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("output_prefix", nargs="?", default="thumbnails", help="Output prefix (default: thumbnails)")
    parser.add_argument("--cols", type=int, default=DEFAULT_COLS, help=f"Columns (default: {DEFAULT_COLS}, max: {MAX_COLS})")
    parser.add_argument("--slides-dir", default=None, help="Save individual slide images to this directory")
    args = parser.parse_args()

    cols = min(args.cols, MAX_COLS)
    input_path = Path(args.input)

    if not input_path.exists() or input_path.suffix.lower() != ".pptx":
        print(f"Error: Invalid PowerPoint file: {args.input}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(f"{args.output_prefix}.jpg")

    slide_info = get_slide_info(input_path)

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        images = convert_to_images(input_path, temp_path)

        if not images:
            print("Error: No slide images generated", file=sys.stderr)
            sys.exit(1)

        # Save individual slide images if requested
        if args.slides_dir:
            slides_dir = Path(args.slides_dir)
            slides_dir.mkdir(parents=True, exist_ok=True)
            for i, img_path in enumerate(images, 1):
                dst = slides_dir / f"slide-{i:02d}.jpg"
                shutil.copy2(img_path, dst)
                print(f"Slide: {dst}")

        # Pair images with slide names
        slides = list(zip(images, [s["name"] for s in slide_info[:len(images)]]))

        # Split into grids if many slides
        max_per_grid = cols * (cols + 1)
        for chunk_idx in range(0, len(slides), max_per_grid):
            chunk = slides[chunk_idx:chunk_idx + max_per_grid]
            grid = create_grid(chunk, cols, THUMBNAIL_WIDTH)

            if len(slides) <= max_per_grid:
                out_file = output_path
            else:
                out_file = output_path.parent / f"{output_path.stem}-{chunk_idx // max_per_grid + 1}{output_path.suffix}"

            grid.save(str(out_file), quality=JPEG_QUALITY)
            print(f"Created: {out_file}")


if __name__ == "__main__":
    main()
